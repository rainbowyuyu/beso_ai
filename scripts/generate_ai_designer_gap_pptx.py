#!/usr/bin/env python3
"""Generate detailed Chinese PPTX: paper claims vs beso_ai implementation gaps."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
OUT_DIR = REPO / "docs" / "gap_analysis_output"
OUT_PPTX = OUT_DIR / "AI_Designer_未实现缺口对照.pptx"
QA_MD = OUT_DIR / "qa_report.md"

# Visual system — engineering teal / slate (avoid purple-AI default)
C_BG = RGBColor(0xF7, 0xF8, 0xFA)
C_INK = RGBColor(0x1A, 0x1F, 0x2E)
C_MUTED = RGBColor(0x5C, 0x65, 0x75)
C_ACCENT = RGBColor(0x00, 0x7A, 0x6C)  # teal
C_WARN = RGBColor(0xB4, 0x53, 0x09)
C_BAD = RGBColor(0xB9, 0x1C, 0x1C)
C_OK = RGBColor(0x04, 0x78, 0x57)
C_LINE = RGBColor(0xD1, 0xD5, 0xDB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL = RGBColor(0xEF, 0xF6, 0xF4)
C_PANEL_BAD = RGBColor(0xFE, 0xF2, 0xF2)
C_PANEL_WARN = RGBColor(0xFF, 0xF7, 0xED)

W, H = Inches(13.333), Inches(7.5)


def _set_run(run, *, size=18, bold=False, color=C_INK, font="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BG
    shape.line.fill.background()


def _bar(slide, y=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ACCENT
    shape.line.fill.background()


def _footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(10), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"对照论文 AI designer-r2 · beso_ai 实现缺口  ·  {page}/{total}"
    _set_run(run, size=11, color=C_MUTED)
    box2 = slide.shapes.add_textbox(Inches(11.2), Inches(7.1), Inches(1.5), Inches(0.3))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "内部缺口分析"
    _set_run(run2, size=11, color=C_MUTED)


def _title(slide, text: str, *, y=Inches(0.25)):
    box = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.2), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=28, bold=True, color=C_INK)


def _subtitle(slide, text: str, *, y=Inches(0.8)):
    box = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.2), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=14, color=C_MUTED)


def _bullets(slide, items: list[str], *, left, top, width, height, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        _set_run(run, size=size, color=C_INK)


def _panel(slide, left, top, width, height, fill=C_PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C_LINE
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.08
    return shape


def _panel_title(slide, text, left, top, width, *, color=C_ACCENT, size=14):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=True, color=color)


def _note(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slides_meta: list[str] = []

    def add_blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    total = 16  # planned

    # ---- 1 Title ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), W, Inches(2.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x36)
    shape.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.45), Inches(12), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "The AI Designer：论文宣称能力 vs 仓库未实现缺口"
    _set_run(run, size=30, bold=True, color=C_WHITE)
    box2 = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(0.9))
    p2 = box2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "对照稿：AI designer-r2 (2026-07-14)  ·  代码库：beso_ai\n焦点：闭环探索、Zwind 尺寸优化、终止门、审计、出图 — 而非已完成的 Validation UI"
    _set_run(run2, size=15, color=RGBColor(0xC8, 0xE6, 0xE0))
    box3 = s.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(0.6))
    p3 = box3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "用途：内部实现差距盘点 / 路线图输入  ·  方法类缺口分析 PPT"
    _set_run(run3, size=13, color=C_MUTED)
    _note(s, "开场：说明本PPT不是论文汇报，而是对照 r2 文稿与 beso_ai 代码，列出「论文写了但工程还没闭环」的项。")
    slides_meta.append("1 标题")

    # ---- 2 Purpose ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "本材料回答什么问题")
    _subtitle(s, "区分三类：已实现 · 部分实现 · 论文宣称但未编码（另附论文自述越界）")
    _panel(s, Inches(0.55), Inches(1.4), Inches(3.9), Inches(4.8), C_PANEL)
    _panel_title(s, "已较齐（对照基准）", Inches(0.75), Inches(1.55), Inches(3.5), color=C_OK)
    _bullets(
        s,
        [
            "OC4→FreeCAD/Gmsh→CalculiX–BESO",
            "AI Review 五维评分 + 机队对照",
            "Validation 报告 / Nature 图 / Word",
            "参数化重建与钢量缩放（解析）",
            "WebSocket 作业日志、INP 存档",
        ],
        left=Inches(0.75),
        top=Inches(2.1),
        width=Inches(3.5),
        height=Inches(3.8),
        size=14,
    )
    _panel(s, Inches(4.7), Inches(1.4), Inches(3.9), Inches(4.8), C_PANEL_WARN)
    _panel_title(s, "部分实现（骨架有断点）", Inches(4.9), Inches(1.55), Inches(3.5), color=C_WARN)
    _bullets(
        s,
        [
            "NL→任务：仅 BESO 超参，非海况说明书",
            "尺度 SLSQP：全局 s，非 (D,t) PSO",
            "Zwind：仅 envelope 导入占位",
            "代理重规划：OC4 工具环，无 F_p",
            "静力 PINN：验证/可选，非优化核",
        ],
        left=Inches(4.9),
        top=Inches(2.1),
        width=Inches(3.5),
        height=Inches(3.8),
        size=14,
    )
    _panel(s, Inches(8.85), Inches(1.4), Inches(3.9), Inches(4.8), C_PANEL_BAD)
    _panel_title(s, "关键未实现（本PPT重点）", Inches(9.05), Inches(1.55), Inches(3.5), color=C_BAD)
    _bullets(
        s,
        [
            "S≥85 作为探索终止门",
            "PSO/SQP(D,t)+PVW+Zwind 约束",
            "多候选生成→最高分选优",
            "失败驱动 replan(θ,Fₚ)",
            "SHA-256 完整性清单",
            "认证级工程图自动出图",
        ],
        left=Inches(9.05),
        top=Inches(2.1),
        width=Inches(3.5),
        height=Inches(3.8),
        size=14,
    )
    _footer(s, 2, total)
    _note(s, "强调：Validation UI 不等于论文里的「探索终止门」。听众容易混淆。")
    slides_meta.append("2 三类对照")

    # ---- 3 Paper loop ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "论文宣称的验证闭环（Fig.1b 四相位）")
    _subtitle(s, "整理自 AI designer-r2 · Phase I–IV；虚线 AIP 在环外作金标准校准")
    phases = [
        ("I", "需求形式化", "NL → 场址海况\n入级条款 → J"),
        ("II", "闭环优化", "BESO→升尺度→\nSQP/PSO→Zwind"),
        ("III", "出图写报告", "工程图 + 设计\n报告包"),
        ("IV", "Reviewer 闸", "S≥85 且子分≥60\n才停探索"),
    ]
    x0 = 0.55
    for i, (num, title, body) in enumerate(phases):
        left = Inches(x0 + i * 3.1)
        _panel(s, left, Inches(1.5), Inches(2.85), Inches(3.2), C_WHITE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.05), Inches(1.7), Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()
        tb = s.shapes.add_textbox(left + Inches(1.05), Inches(1.82), Inches(0.7), Inches(0.5))
        tp = tb.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        r = tp.add_run()
        r.text = num
        _set_run(r, size=18, bold=True, color=C_WHITE)
        _panel_title(s, title, left + Inches(0.15), Inches(2.55), Inches(2.55), size=16)
        box = s.shapes.add_textbox(left + Inches(0.2), Inches(3.05), Inches(2.5), Inches(1.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body
        _set_run(run, size=13, color=C_MUTED)
        if i < 3:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + Inches(2.9),
                Inches(2.8),
                Inches(0.22),
                Inches(0.28),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_LINE
            arr.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "论文终止条件（Methods）：ρₚ≡0（各相无重试）∧ W(x) 相对下降 <10⁻⁴（连续三步）"
        "∧ S≥85（子分均≥60）。CCS AIP 是一次性外校准，不是每代目标函数。"
    )
    _set_run(run, size=14, color=C_INK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = "仓库现状：II 前半段 + IV 评分 UI 较完整；II 尺寸/Zwind、IV 终止门、I 场址说明、III 认证出图是最大断点。"
    _set_run(run2, size=14, bold=True, color=C_BAD)
    _footer(s, 3, total)
    _note(s, "对照 Fig.1b 讲清：日常环内是 Automated Reviewer；AIP 在环外。")
    slides_meta.append("3 四相位闭环")

    # ---- 4 Status matrix ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "总览矩阵：论文项 → 代码判定")
    _subtitle(s, "绿=已实现  黄=部分  红=未实现（相对论文宣称，非论文 Limitations 段落）")
    rows = [
        ("能力模块", "论文主张", "仓库判定"),
        ("Phase I 场址海况+入级→J", "结构化说明书", "黄/红"),
        ("OC4–Gmsh–CalculiX–BESO", "拓扑发现", "绿"),
        ("20 MW 参数化升尺度 STEP", "拓扑→实体闭环", "黄"),
        ("PSO/SQP (D,t) + PVW", "尺寸优核心", "红"),
        ("Zwind 时域耦合", "极限态真值", "红（适配器黄）"),
        ("失败 replan(θ,Fₚ)", "自治恢复", "黄/红"),
        ("认证工程图自动出图", "Phase III", "红"),
        ("AI Review 五维打分", "内审评分", "绿"),
        ("S≥85 停探索", "终止门", "红"),
        ("多候选选最优→AIP 包", "探索+归档", "红"),
        ("SHA-256 完整性清单", "审计", "红"),
    ]
    table = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.3)).table
    table.columns[0].width = Inches(4.2)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(3.8)
    for r_i, row in enumerate(rows):
        for c_i, cell_text in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if r_i == 0:
                        _set_run(run, size=13, bold=True, color=C_WHITE)
                    else:
                        color = C_INK
                        if c_i == 2:
                            if "绿" in cell_text:
                                color = C_OK
                            elif "黄" in cell_text:
                                color = C_WARN
                            elif "红" in cell_text:
                                color = C_BAD
                        _set_run(run, size=12, bold=(c_i == 2), color=color)
            if r_i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x36)
            elif r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF3)
    _footer(s, 4, total)
    slides_meta.append("4 总览矩阵")

    # ---- 5 Phase I ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "Phase I 缺口：需求形式化不完整")
    _subtitle(s, "论文：owner 意图 + 场址 sea-state + 船级社条款 → job descriptor J")
    _panel(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(4.9), C_WHITE)
    _panel_title(s, "已有", Inches(0.75), Inches(1.55), Inches(5.5), color=C_OK)
    _bullets(
        s,
        [
            "agent.decide_params：NL → BESO 超参",
            "oc4_nl_loads：NL → CalculiX 载荷片段",
            "validation_rules / dnv_clause_index：打分用条款库",
        ],
        left=Inches(0.75),
        top=Inches(2.1),
        width=Inches(5.5),
        height=Inches(3.8),
        size=15,
    )
    _panel(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(4.9), C_PANEL_BAD)
    _panel_title(s, "未实现 / 断点", Inches(7.0), Inches(1.55), Inches(5.5), color=C_BAD)
    _bullets(
        s,
        [
            "无统一 J={phase, θ, retry_policy} 对象",
            "无 Hs/Tp/风向错位等场址海况写入任务",
            "入级条款未成为 Phase I 产出，只在 IV 评分消费",
            "无 schema 校验后的跨相位任务卡流转",
        ],
        left=Inches(7.0),
        top=Inches(2.1),
        width=Inches(5.5),
        height=Inches(3.8),
        size=15,
    )
    _footer(s, 5, total)
    _note(s, "Phase I 是「说明书→可调度任务」；现状更像「NL 微改 BESO 旋钮」。")
    slides_meta.append("5 Phase I")

    # ---- 6 Phase II overview ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "Phase II：前半齐整，后半断裂")
    _subtitle(s, "拓扑到网格可跑；尺寸优化–Zwind–重规划是论文最重、仓库最空的一段")
    items = [
        ("① OC4 先验 / 设计域", "绿", "oc4_design_domain_* / INP elset"),
        ("② FreeCAD 非设计域", "绿", "FCStd→INP 链路"),
        ("③ Gmsh 网格", "绿", "gmsh_iges_to_inp"),
        ("④ CalculiX–BESO", "绿", "beso/ + jobs"),
        ("⑤ 升尺度 STEP/FCStd", "黄", "mixed_platform_steel + beso7；无自治闭环"),
        ("⑥ SQP/PSO (D,t)", "红", "仅有全局 scale 的 SLSQP+pitch"),
        ("⑦ Zwind 时域极限态", "红", "zwind_adapter 只导入 envelope"),
        ("⑧ 失败自治重规划", "黄", "工具失败可再试；无 replan(θ,Fₚ)"),
    ]
    for i, (title, tag, detail) in enumerate(items):
        col = i % 4
        row = i // 4
        left = Inches(0.55 + col * 3.15)
        top = Inches(1.4 + row * 2.6)
        fill = C_PANEL if tag == "绿" else (C_PANEL_WARN if tag == "黄" else C_PANEL_BAD)
        color = C_OK if tag == "绿" else (C_WARN if tag == "黄" else C_BAD)
        _panel(s, left, top, Inches(3.0), Inches(2.3), fill)
        _panel_title(s, f"{title}  [{tag}]", left + Inches(0.15), top + Inches(0.15), Inches(2.7), color=color, size=13)
        box = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.7), Inches(2.7), Inches(1.3))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = detail
        _set_run(run, size=13, color=C_INK)
    _footer(s, 6, total)
    slides_meta.append("6 Phase II 总览")

    # ---- 7 Size + PVW + Zwind ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "关键缺口 A：尺寸优化 × 虚功 × Zwind")
    _subtitle(s, "Methods「Size optimization / Zwind」——论文说 PSO+PVW 规避每粒子全时域，仓库尚未编码")
    _panel(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(1.35), C_PANEL_BAD)
    box = s.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "论文：min W(x)，x={壳厚 t, 外径 D}；约束 θ_mean≤5°、θ_max≤10°、UC≤1、f 避 1P/3P、D≤1/DFF；"
        "PSO 迭代；PVW 单次静力得灵敏度；Zwind 做极端/工况校核。"
    )
    _set_run(run, size=14, color=C_BAD)
    _panel(s, Inches(0.55), Inches(2.95), Inches(6.0), Inches(3.5), C_WHITE)
    _panel_title(s, "仓库实际", Inches(0.75), Inches(3.1), Inches(5.5), color=C_WARN)
    _bullets(
        s,
        [
            "mixed_platform_steel.optimize_scale：只优化全局水平缩放 x",
            "约束仅为解析静倾 ≤5°（非 Zwind pitch）",
            "全库无 PSO；无 (D,t) 设计向量优化器",
            "无 PVW 单位载荷解析梯度路径",
            "zwind_adapter：JSON/YAML 导入八指标，不调求解器",
        ],
        left=Inches(0.75),
        top=Inches(3.6),
        width=Inches(5.5),
        height=Inches(2.6),
        size=14,
    )
    _panel(s, Inches(6.8), Inches(2.95), Inches(5.9), Inches(3.5), C_PANEL)
    _panel_title(s, "落地建议（优化环，非验证页）", Inches(7.0), Inches(3.1), Inches(5.5), color=C_ACCENT)
    _bullets(
        s,
        [
            "离线：DOE(D,t,s) → 批量 Zwind 包络标签",
            "训练：多输出 PINN + PVW/频率锚（可扩展 surrogate）",
            "在线：PSO/SQP 调 predict(x)，每 N 代真 Zwind 校核",
            "接 optimize_with_surrogate，与 run_validation 解耦",
        ],
        left=Inches(7.0),
        top=Inches(3.6),
        width=Inches(5.5),
        height=Inches(2.6),
        size=14,
    )
    _footer(s, 7, total)
    _note(s, "用户已明确：PINN/PVW 应用在优化过程，不是 validation UI。")
    slides_meta.append("7 尺寸+PVW+Zwind")

    # ---- 8 Replanning ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "关键缺口 B：失败驱动重规划未形式化")
    _subtitle(s, "论文 Table S.1：Gmsh 翻转单元 / CalculiX 残差平台 / Zwind pitch 超限 → 自动改 θ 再跑")
    rows = [
        ("案例", "论文诊断信号", "仓库现状"),
        ("网格失败", "mesh_quality_min<τ", "可有工具失败摘要；无质量阈值→加密策略库"),
        ("静力不收敛", "residual_norm 平台", "BESO/CalculiX 作业可重跑；无自动改步长/初值策略"),
        ("Zwind 中止", "pitch_max>限值", "无 Zwind 子进程 → 案例整条缺失"),
        ("反馈元组 Fₚ", "Lₚ,Mₚ,ρₚ", "无类型/无 API；无 replan(θ,Fₚ)"),
        ("阈值 τ", "船级社+求解器文档", "散落在 validation/surrogate；未驱动编排"),
    ]
    table = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.4), Inches(12.2), Inches(5.0)).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(4.6)
    table.columns[2].width = Inches(5.2)
    for r_i, row in enumerate(rows):
        for c_i, text in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_run(run, size=12, bold=(r_i == 0), color=C_WHITE if r_i == 0 else C_INK)
            if r_i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x36)
    _footer(s, 8, total)
    slides_meta.append("8 重规划")

    # ---- 9 Phase III drawings ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "Phase III 缺口：认证工程图自动化")
    _subtitle(s, "报告管线已强；GA/平立剖/标注出图包仍缺（论文 Discussion 也承认细部制造不在范围）")
    _panel(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(4.9), C_PANEL)
    _panel_title(s, "已实现（勿低估）", Inches(0.75), Inches(1.55), Inches(5.5), color=C_OK)
    _bullets(
        s,
        [
            "validation pipeline：MD / JSON / 双轨 Word",
            "Nature 风格图：基准、雷达、有效性、PINN 面板",
            "详细设计基础风格 Word（封面目录四章）",
            "CAD 设计台 / text-to-cad（STEP 生成，非审图图纸）",
        ],
        left=Inches(0.75),
        top=Inches(2.1),
        width=Inches(5.5),
        height=Inches(3.8),
        size=15,
    )
    _panel(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(4.9), C_PANEL_BAD)
    _panel_title(s, "相对论文仍缺", Inches(7.0), Inches(1.55), Inches(5.5), color=C_BAD)
    _bullets(
        s,
        [
            "总布置图、关键系统原理图自动投影",
            "平/立/剖成套标注、图框与标题栏标准包",
            "AIP 文档包一键组装（图纸+计算书清单）",
            "注意：焊图/NDE/吊装属论文 Limitations，勿当「缺实现」硬伤",
        ],
        left=Inches(7.0),
        top=Inches(2.1),
        width=Inches(5.5),
        height=Inches(3.8),
        size=15,
    )
    _footer(s, 9, total)
    slides_meta.append("9 Phase III")

    # ---- 10 Phase IV gate ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "关键缺口 C：S≥85 尚未成为探索终止门")
    _subtitle(s, "打分体系已校准；缺少「生成环」读取分数并停搜/归档")
    # two columns big claim
    _panel(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(2.2), C_PANEL)
    _panel_title(s, "已实现", Inches(0.75), Inches(1.55), Inches(5.5), color=C_OK)
    _bullets(
        s,
        [
            "五维 AI Review + 法规对照通道",
            "等级 A/B/C/D（≥85→A）",
            "机队 Spearman≈0.72 有效性统计可复现",
        ],
        left=Inches(0.75),
        top=Inches(2.05),
        width=Inches(5.5),
        height=Inches(1.4),
        size=14,
    )
    _panel(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.2), C_PANEL_BAD)
    _panel_title(s, "未实现（叙事关键）", Inches(7.0), Inches(1.55), Inches(5.5), color=C_BAD)
    _bullets(
        s,
        [
            "编排器不读 overall_score 决定停探索",
            "无「任一子分 <60 → 拒绝放行」闸",
            "Validation 是独立 API/UI，非环内闸",
        ],
        left=Inches(7.0),
        top=Inches(2.05),
        width=Inches(5.5),
        height=Inches(1.4),
        size=14,
    )
    _panel(s, Inches(0.55), Inches(3.9), Inches(12.2), Inches(2.4), C_WHITE)
    _panel_title(s, "论文 Methods 三条件 vs 代码", Inches(0.75), Inches(4.05), Inches(11.8), color=C_ACCENT)
    _bullets(
        s,
        [
            "① ρₚ=0 各相无 pending 重试 —— 无统一相位状态机",
            "② W(x) 相对降幅 <10⁻⁴ 连续三步 —— 尺寸优环不存在，无从谈起",
            "③ S≥85 且 sᵢ≥60 —— 仅标签与报告；无 halt_and_archive()",
        ],
        left=Inches(0.75),
        top=Inches(4.55),
        width=Inches(11.8),
        height=Inches(1.6),
        size=15,
    )
    _footer(s, 10, total)
    _note(s, "这是最易被外审追问的点：有 Reviewer 分数，但没有「用分数关循环」。")
    slides_meta.append("10 终止门")

    # ---- 11 Orchestrator ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "编排器缺口：有工具环，无验证闭环状态机")
    _subtitle(s, "论文：允许工具集合 + schema + Fₚ + τ → replan；仓库：OC4/助手白名单工具 + 多轮 JSON")
    _bullets(
        s,
        [
            "缺少跨 Phase I–IV 的统一 DAG：后相不得早启（ρₚ≠0 时）",
            "缺少标准化反馈元组 Fₚ=(日志Lₚ, 指标Mₚ, 重试ρₚ)",
            "缺少规则增强阈值 τ 注册表，并把求解器诊断映射到策略库",
            "工具温度/可复现：生产认证建议 multi-seed + 冻结 schema（论文 Limitations；工程侧未固化策略）",
            "结果：系统表现为「强模块 + 弱 orchestrator」，而非「agentic closed loop」",
        ],
        left=Inches(0.7),
        top=Inches(1.5),
        width=Inches(12),
        height=Inches(5.0),
        size=17,
    )
    _footer(s, 11, total)
    slides_meta.append("11 编排器")

    # ---- 12 Audit ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "审计缺口：SHA-256 完整性清单")
    _subtitle(s, "论文强调 AIP 澄清可回溯到具体计算步；代码已有部分底座")
    _panel(s, Inches(0.55), Inches(1.5), Inches(4.0), Inches(4.6), C_PANEL)
    _panel_title(s, "已有", Inches(0.75), Inches(1.7), Inches(3.6), color=C_OK)
    _bullets(
        s,
        ["WebSocket 作业日志", "INP / beso_conf 存档", "validation 产物目录", "任务参数 manifest（路由）"],
        left=Inches(0.75),
        top=Inches(2.3),
        width=Inches(3.6),
        height=Inches(3.4),
        size=15,
    )
    _panel(s, Inches(4.8), Inches(1.5), Inches(4.0), Inches(4.6), C_PANEL_BAD)
    _panel_title(s, "缺失", Inches(5.0), Inches(1.7), Inches(3.6), color=C_BAD)
    _bullets(
        s,
        ["几何/网格/结果 SHA-256 清单", "不可篡改内容指纹链", "澄清项→计算步自动溯源 UI", "冻结依赖版本锁定策略"],
        left=Inches(5.0),
        top=Inches(2.3),
        width=Inches(3.6),
        height=Inches(3.4),
        size=15,
    )
    _panel(s, Inches(9.05), Inches(1.5), Inches(3.7), Inches(4.6), C_PANEL_WARN)
    _panel_title(s, "优先级", Inches(9.25), Inches(1.7), Inches(3.3), color=C_WARN)
    box = s.shapes.add_textbox(Inches(9.25), Inches(2.4), Inches(3.3), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "中高：相对 Zwind/PSO 实现成本低，但对「可审性」故事强相关；可先做 hash manifest 生成器。"
    _set_run(run, size=14, color=C_INK)
    _footer(s, 12, total)
    slides_meta.append("12 审计")

    # ---- 13 Multi-candidate ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "多候选探索与 AIP 选交：未实现")
    _subtitle(s, "论文叙事：千候选廉价内审 → 单巅峰交 CCS；仓库「candidate」多为单方案机队对照标签")
    _bullets(
        s,
        [
            "无设计族（topology/scale/t 分支）并行生成与登记",
            "无按 S 排序选最高分并 halt_and_archive 的策略",
            "无 AIP 技术卷宗清单自动组装（图纸+极限态计算书+系泊+稳性）",
            "无船级社往来澄清回灌到编排器的接口",
            "注意：domain_mapping_candidates = Elset 映射，不是设计候选族",
        ],
        left=Inches(0.7),
        top=Inches(1.5),
        width=Inches(12),
        height=Inches(5.0),
        size=17,
    )
    _footer(s, 13, total)
    slides_meta.append("13 多候选")

    # ---- 14 Paper limitations ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "勿混淆：论文自述边界 ≠ 代码缺口")
    _subtitle(s, "Discussion Limitations —— PPT/答辩应分开列，避免「本来就没宣称」被当成欠债")
    rows = [
        ("论文自述越界 / 未来工作", "建议姿态"),
        ("碎波抨击、强非线性粘性阻力 → 需水池/CFD", "标 future；勿写进「缺编码」主清单"),
        ("焊图、NDE、船坞吊装、静转动安装", "标 OOS；detailed design 人类负责"),
        ("LLM 随机性 → multi-seed + 冻结 schema", "工程可补强，属治理而非缺模块"),
        ("图强基线持续前移 → 需周期重对标", "数据/评分维护问题"),
        ("细部设计与最终型式批准在 AIP 之后", "产品边界；已与内部终止门区分"),
    ]
    table = s.shapes.add_table(len(rows), 2, Inches(0.55), Inches(1.4), Inches(12.2), Inches(5.0)).table
    table.columns[0].width = Inches(7.5)
    table.columns[1].width = Inches(4.7)
    for r_i, row in enumerate(rows):
        for c_i, text in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_run(run, size=13, bold=(r_i == 0), color=C_WHITE if r_i == 0 else C_INK)
            if r_i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x36)
    _footer(s, 14, total)
    slides_meta.append("14 论文边界")

    # ---- 15 Roadmap ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "建议落地优先级（对齐论文故事）")
    _subtitle(s, "P0 讲清闭环；P1 把尺寸环跑通；P2 补审计与出图产能")
    roadmap = [
        ("P0", "终止门接线", "halt_and_archive：S≥85 ∧ min(sᵢ)≥60；探索环读分数", C_BAD),
        ("P0", "相位状态机", "Fₚ / ρₚ / DAG；validation 从「页面」升为「环内闸」", C_BAD),
        ("P1", "尺寸优化核", "x={D,t,s}；PSO/SQP；pitch/UC/freq/fatigue 约束接口", C_WARN),
        ("P1", "Zwind+PVW", "真值批量 + PVW/代理加速；末步全时域校核", C_WARN),
        ("P2", "审计清单", "SHA-256 manifest；澄清→计算步溯源", C_ACCENT),
        ("P2", "工程图 MVP", "从 STEP 投影总布置/三视图图框包（非焊图）", C_ACCENT),
        ("P2", "多候选登记", "设计族账号 + 最高分选交", C_ACCENT),
    ]
    for i, (lvl, title, detail, color) in enumerate(roadmap):
        top = Inches(1.35 + i * 0.72)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), top, Inches(0.9), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.55), top + Inches(0.1), Inches(0.9), Inches(0.4))
        tp = tb.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        r = tp.add_run()
        r.text = lvl
        _set_run(r, size=14, bold=True, color=C_WHITE)
        box = s.shapes.add_textbox(Inches(1.65), top, Inches(2.6), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_run(run, size=16, bold=True, color=C_INK)
        box2 = s.shapes.add_textbox(Inches(4.3), top, Inches(8.4), Inches(0.55))
        p2 = box2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = detail
        _set_run(run2, size=14, color=C_MUTED)
    _footer(s, 15, total)
    slides_meta.append("15 路线图")

    # ---- 16 Summary ----
    s = add_blank()
    _add_bg(s)
    _bar(s)
    _title(s, "一句话结论")
    _panel(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.2), RGBColor(0x0B, 0x3D, 0x36))
    box = s.shapes.add_textbox(Inches(0.85), Inches(1.85), Inches(11.6), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "beso_ai 已具备「拓扑发现 + 内审评分 + 报告出图」的强模块；"
        "相对 AI designer-r2，「用分数终止探索的编排闭环」与「Zwind/PVW 驱动的尺寸优化」仍是最大未兑现主张。"
    )
    _set_run(run, size=20, bold=True, color=C_WHITE)
    _bullets(
        s,
        [
            "不要用 Validation UI 替代论文里的 Phase IV 终止门叙事——需接线到生成环",
            "PINN/虚功应服务 PSO/SQP 在线评价，而非只做事后验证页可视化",
            "论文 Limitations（水池/焊图）与「宣称已做却未编码」必须分列，减少审稿误会",
        ],
        left=Inches(0.7),
        top=Inches(4.1),
        width=Inches(12),
        height=Inches(2.4),
        size=16,
    )
    _footer(s, 16, total)
    _note(s, "收束：建议下一迭代以「终止门接线 + Zwind 标签尺寸优」为双 P0。")
    slides_meta.append("16 结论")

    prs.save(OUT_PPTX)

    QA_MD.write_text(
        "\n".join(
            [
                "# QA Report — AI Designer 缺口对照 PPT",
                "",
                f"- PPTX: `{OUT_PPTX.as_posix()}`",
                f"- Slides: {len(prs.slides)}",
                f"- Paper type: methods (gap-analysis adaptation)",
                f"- Source: AI designer-r2(1).docx (2026-07-14)",
                f"- Codebase: beso_ai (gap map via repo explore)",
                "",
                "## Slide list",
                *[f"{i+1}. {m}" for i, m in enumerate(slides_meta)],
                "",
                "## Design notes",
                "- No paper figure crops embedded (gap analysis; native tables/panels).",
                "- Palette: engineering teal/slate; avoided purple-gradient AI template look.",
                "- Distinguishes paper Limitations from implementation gaps.",
                "",
                "## Self-review",
                "- High: none expected for fabricated metrics — numbers tied to paper (S≥85, ρ=0.717) or qualitative gap labels.",
                "- Medium: long table cells may wrap tightly on slide 4/8 — kept font ≥12pt.",
                "- Verification: file written; reopen with python-pptx recommended.",
                "",
            ]
        ),
        encoding="utf-8",
    )
    return OUT_PPTX


if __name__ == "__main__":
    path = build()
    print(path)
    # reopen check
    from pptx import Presentation as P

    prs = P(str(path))
    print("slides", len(prs.slides), "ok")
